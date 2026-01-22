#!/usr/bin/env python3
"""
プレゼン資料生成スクリプト
PowerPoint形式（.pptx）のプレゼン資料を生成します。
"""

import os
import sys
from pathlib import Path
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# プロジェクトルートのパス
PROJECT_ROOT = Path(__file__).parent.parent
RESULTS_DIR = PROJECT_ROOT / "results"
FIGURES_DIR = PROJECT_ROOT / "outputs" / "figures"
OUTPUT_DIR = PROJECT_ROOT / "outputs" / "presentation"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 評価結果データの読み込み
EVAL_CSV = RESULTS_DIR / "evaluation_summary.csv"

# グラフファイル
MCD_FIGURE = FIGURES_DIR / "mcd_comparison.png"
F0_FIGURE = FIGURES_DIR / "f0_comparison.png"
COVERAGE_FIGURE = FIGURES_DIR / "coverage_vs_mcd.png"
DATASIZE_FIGURE = FIGURES_DIR / "datasize_vs_mcd.png"


def load_evaluation_data():
    """評価結果データを読み込む"""
    df = pd.read_csv(EVAL_CSV)
    # 条件名を日本語に変換
    condition_map = {
        "train_80sent": "80文コーパス",
        "train_4sent_37phonemes": "37音素4文",
        "train_4sent_random": "ランダム4文",
        "train_10sent_top": "上位10文"
    }
    df["condition_jp"] = df["condition"].map(condition_map)
    return df


def add_title_slide(prs):
    """タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "少量音素コーパスを用いた\n日本語TTS構築実験"
    subtitle.text = "音素カバレッジとコーパス設計が品質に与える影響"
    
    # フォントサイズ調整
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # 日付を追加（サブタイトルの下）
    date_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_text.text_frame
    date_frame.text = f"作成日: {datetime.now().strftime('%Y年%m月%d日')}"
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def add_background_slide(prs):
    """背景・研究目的スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    
    title = slide.shapes.title
    title.text = "背景・研究目的"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "研究背景"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 少量データでのTTS fine-tuningの重要性"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 音素カバレッジとデータ量の最適なバランスを探索"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "研究目的"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "音素カバレッジとコーパス設計（4文 vs 80文など）が、"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "fine-tuningされた日本語TTSモデルの出力品質に与える影響を、"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "客観指標（MCD, log-F0 RMSE）を用いて明らかにする"
    p.font.size = Pt(18)
    p.level = 1


def add_objective_slide_final(prs):
    """最終発表用目的スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "目的"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "【メインゴール】"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 少量データで効率的にTTSモデルを構築するための必要最低限コーパスを設計する"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 音素カバレッジに着目したデータ量削減の有効性を検証"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Fine-tuningによる学習時間短縮の実現"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "【具体的達成目標】"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• コーパス設計：JVS parallel100（100文）から必要最小限の文セット抽出"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 音素カバレッジの定量化：貪欲法により4文で全37音素をカバー"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 削減効果の検証：「100文vs削減コーパス」で品質比較（MCD、log-F0 RMSE）"
    p.font.size = Pt(18)
    p.level = 1


def add_experiment_design_slide(prs):
    """実験設計スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "方法"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # 表を作成
    rows = 5
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # ヘッダー行
    table.cell(0, 0).text = "条件"
    table.cell(0, 1).text = "文数"
    table.cell(0, 2).text = "特徴"
    
    # データ行
    table.cell(1, 0).text = "E1: 80文コーパス"
    table.cell(1, 1).text = "80"
    table.cell(1, 2).text = "データ量最大、分布的にバランス良い"
    
    table.cell(2, 0).text = "E2: 37音素4文"
    table.cell(2, 1).text = "4"
    table.cell(2, 2).text = "音素カバレッジ最大、データ量最小"
    
    table.cell(3, 0).text = "E3: ランダム4文"
    table.cell(3, 1).text = "4"
    table.cell(3, 2).text = "カバレッジ・分布ともに無作為（対照群）"
    
    table.cell(4, 0).text = "E4: 上位10文"
    table.cell(4, 1).text = "10"
    table.cell(4, 2).text = "音素特徴量上位10文（情報量重視）"
    
    # フォントサイズ設定
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if row == 0:  # ヘッダー
                cell.text_frame.paragraphs[0].font.bold = True
    
    # 方法の詳細を追加
    left_box = Inches(1)
    top_box = Inches(5.5)
    width_box = Inches(8)
    height_box = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left_box, top_box, width_box, height_box)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "データ: JVS Corpus jvs002話者 parallel100（100文）"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "音素解析: pyopenjtalkによる音素抽出、37音素インベントリの確認"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "Fine-tuning: ESPnet2-TTS (Tacotron2), 10エポック"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "評価: MCD（音色類似度）、log-F0 RMSE（韻律誤差）、テストセット18文"
    p.font.size = Pt(14)


def add_evaluation_metrics_slide(prs):
    """評価指標スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "評価指標"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MCD (Mel-Cepstral Distortion)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• スペクトル類似度を測定（低いほど良い）"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 単位: dB"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "log-F0 RMSE"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• ピッチ輪郭（基本周波数）の誤差を測定（低いほど良い）"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• DTWアライメントを使用"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "テストセット: 18文（全条件共通）"
    p.font.size = Pt(18)
    p.font.italic = True
    p.level = 0


def add_mcd_results_slide(prs, df):
    """MCD結果スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "実験結果 - MCD"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # 数値表を作成
    rows = 5
    cols = 2
    left = Inches(1)
    top = Inches(2)
    width = Inches(3.5)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.cell(0, 0).text = "条件"
    table.cell(0, 1).text = "MCD (dB)"
    
    for i, row in df.iterrows():
        table.cell(i + 1, 0).text = row["condition_jp"]
        table.cell(i + 1, 1).text = f"{row['mcd']:.3f}"
    
    # フォントサイズ設定
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if row == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    # グラフを挿入
    if MCD_FIGURE.exists():
        left_img = Inches(5)
        top_img = Inches(2)
        width_img = Inches(4.5)
        height_img = Inches(3.5)
        slide.shapes.add_picture(str(MCD_FIGURE), left_img, top_img, width_img, height_img)
    
    # 解釈テキスト
    left_text = Inches(1)
    top_text = Inches(4.8)
    width_text = Inches(8.5)
    height_text = Inches(0.8)
    
    text_box = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
    tf = text_box.text_frame
    tf.text = "解釈: データ量が多いほど品質が向上（80文: 5.049 dB < 4文: 5.227-5.266 dB）"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True


def add_f0_results_slide(prs, df):
    """log-F0 RMSE結果スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "実験結果 - log-F0 RMSE"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # 数値表を作成
    rows = 5
    cols = 2
    left = Inches(1)
    top = Inches(2)
    width = Inches(3.5)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.cell(0, 0).text = "条件"
    table.cell(0, 1).text = "log-F0 RMSE"
    
    for i, row in df.iterrows():
        table.cell(i + 1, 0).text = row["condition_jp"]
        table.cell(i + 1, 1).text = f"{row['log_f0_rmse']:.4f}"
    
    # フォントサイズ設定
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if row == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    # グラフを挿入
    if F0_FIGURE.exists():
        left_img = Inches(5)
        top_img = Inches(2)
        width_img = Inches(4.5)
        height_img = Inches(3.5)
        slide.shapes.add_picture(str(F0_FIGURE), left_img, top_img, width_img, height_img)
    
    # 解釈テキスト
    left_text = Inches(1)
    top_text = Inches(4.8)
    width_text = Inches(8.5)
    height_text = Inches(0.8)
    
    text_box = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
    tf = text_box.text_frame
    tf.text = "解釈: 10エポックでは中間的な音程に収束（完全な話者適応には不十分）"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.italic = True


def add_visualization_slide(prs):
    """可視化スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # タイトル
    left_title = Inches(0.5)
    top_title = Inches(0.3)
    width_title = Inches(9)
    height_title = Inches(0.6)
    
    title_box = slide.shapes.add_textbox(left_title, top_title, width_title, height_title)
    tf_title = title_box.text_frame
    tf_title.text = "実験結果 - 可視化"
    tf_title.paragraphs[0].font.size = Pt(36)
    tf_title.paragraphs[0].font.bold = True
    
    # 左側のグラフ（音素カバレッジとMCD）
    if COVERAGE_FIGURE.exists():
        left_img1 = Inches(0.5)
        top_img1 = Inches(1.2)
        width_img1 = Inches(4.5)
        height_img1 = Inches(3.5)
        slide.shapes.add_picture(str(COVERAGE_FIGURE), left_img1, top_img1, width_img1, height_img1)
    
    # 右側のグラフ（データ量とMCD）
    if DATASIZE_FIGURE.exists():
        left_img2 = Inches(5.5)
        top_img2 = Inches(1.2)
        width_img2 = Inches(4.5)
        height_img2 = Inches(3.5)
        slide.shapes.add_picture(str(DATASIZE_FIGURE), left_img2, top_img2, width_img2, height_img2)


def add_discussion_slide_1(prs):
    """考察スライド1を追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "考察（1）"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "10エポックの限界"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• log-F0 RMSEの結果から、10エポックでは完全な話者適応には不十分"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• ベースラインモデル（高音）とターゲット話者（低音）の中間的な音程に収束"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• 完全な話者適応には、20-50エポック以上の学習が必要と推測"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "音素カバレッジとデータ量の関係"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 音素カバレッジ100%でも、データ量が少ない（4文）場合は品質が低下"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 37音素4文（MCD=5.227 dB）とランダム4文（MCD=5.266 dB）の差は小さい（約0.04 dB）"
    p.font.size = Pt(16)
    p.level = 2


def add_discussion_slide_2(prs):
    """考察スライド2を追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "考察（続き）"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "データ量と品質のバランス"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 80文と10文の差は小さい（約0.05 dB）→ 10文程度でも十分な品質が得られる可能性"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "10エポックの限界"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 完全な話者適応には不十分（中間的な音程に収束）"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• より多くのエポック数（20-50エポック）が必要と推測"
    p.font.size = Pt(16)
    p.level = 2


def add_conclusion_slide(prs):
    """結論スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "結論"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "主要な発見"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "1. 10エポックでは完全な話者適応には不十分"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   log-F0 RMSEの結果から、音程が中間的になることが数値で裏付けられた"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "2. 音素カバレッジ100%でもデータ量が少ない（4文）場合は品質が低下"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   データ量が多い（80文）場合は、音素カバレッジが100%でなくても品質が高い"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "3. 10文程度でも十分な品質が得られる可能性"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "   80文と10文の差は小さい（約0.05 dB）ことから、実用的な最小データ量として有望"
    p.font.size = Pt(16)
    p.level = 2


def add_future_work_slide(prs):
    """今後の課題スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "今後の課題"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "エポック数の増加"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 20-50エポック以上の学習を実施し、完全な話者適応が可能かを検証"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• log-F0 RMSEが0.1以下になるかを確認"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "主観評価の実施"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• MOS（Mean Opinion Score）評価を実施し、客観評価と主観評価の関係を検証"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 音声の自然さ、話者類似度などの主観的評価を実施"
    p.font.size = Pt(16)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "統計的有意性検定"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 条件間の差が統計的に有意かどうかを検証"
    p.font.size = Pt(18)
    p.level = 1


def main():
    """メイン関数"""
    print("プレゼン資料を生成中...")
    
    # 評価データの読み込み
    df = load_evaluation_data()
    
    # プレゼンテーションオブジェクトの作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # スライドを追加
    print("  タイトルスライドを追加...")
    add_title_slide(prs)
    
    print("  背景・研究目的スライドを追加...")
    add_background_slide(prs)
    
    print("  実験設計スライドを追加...")
    add_experiment_design_slide(prs)
    
    print("  評価指標スライドを追加...")
    add_evaluation_metrics_slide(prs)
    
    print("  MCD結果スライドを追加...")
    add_mcd_results_slide(prs, df)
    
    print("  log-F0 RMSE結果スライドを追加...")
    add_f0_results_slide(prs, df)
    
    print("  可視化スライドを追加...")
    add_visualization_slide(prs)
    
    print("  考察スライド1を追加...")
    add_discussion_slide_1(prs)
    
    print("  考察スライド2を追加...")
    add_discussion_slide_2(prs)
    
    print("  結論スライドを追加...")
    add_conclusion_slide(prs)
    
    print("  今後の課題スライドを追加...")
    add_future_work_slide(prs)
    
    # ファイルを保存
    output_file = OUTPUT_DIR / "slides.pptx"
    prs.save(str(output_file))
    
    print(f"\nプレゼン資料を生成しました: {output_file}")
    print(f"総スライド数: {len(prs.slides)}枚")


# ============================================================================
# 最終発表資料用の関数
# ============================================================================

def load_phoneme_data():
    """音素分布データを読み込む"""
    import json
    with open(RESULTS_DIR / "phoneme_distribution.json", "r", encoding="utf-8") as f:
        return json.load(f)


def load_corpus_selection():
    """コーパス選定データを読み込む"""
    import json
    with open(RESULTS_DIR / "corpus_selection.json", "r", encoding="utf-8") as f:
        return json.load(f)


def add_title_slide_final(prs):
    """最終発表用タイトルスライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "少量特徴量を用いた\n音声合成モデルの構築"
    subtitle.text = "音素カバレッジに着目したデータ量・学習時間の削減"
    
    # フォントサイズ調整
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # 発表者情報を追加
    author_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.5))
    author_frame = author_text.text_frame
    author_frame.text = "発表者: 37_TK230037_國政蒼矢(くにまさそうや)"
    author_frame.paragraphs[0].font.size = Pt(18)
    
    # 所属・授業名を追加
    affiliation_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.5))
    affiliation_frame = affiliation_text.text_frame
    affiliation_frame.text = "東京国際工科専門職大学 | 人工知能応用"
    affiliation_frame.paragraphs[0].font.size = Pt(16)
    
    # 日付を追加
    date_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    date_frame = date_text.text_frame
    date_frame.text = "発表日: 2025年12月18日"
    date_frame.paragraphs[0].font.size = Pt(14)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


def add_table_of_contents(prs):
    """目次スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
    
    title = slide.shapes.title
    title.text = "目次"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    items = [
        "背景（従来手法の課題）",
        "目的",
        "方法",
        "結果（音素解析、全音素カバーコーパス、Fine-tuning時間、客観評価、音声合成）",
        "考察・今後の課題"
    ]
    
    for i, item in enumerate(items):
        p = tf.paragraphs[i] if i == 0 else tf.add_paragraph()
        p.text = f"{i+1}. {item}"
        p.font.size = Pt(18)
        p.level = 0


def add_background_challenges(prs):
    """従来手法の課題スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "背景：従来手法の課題"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "【課題1】大量データの必要性"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 40分程度のデータでも「少量」と定義される"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 100以上のコーパスで学習させるのが一般的"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• データ収集コストが高い"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "引用: 令和4年3月 井上勝喜「深層学習に基づく感情音声合成のための少量データを用いた学習方式の研究」"
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "【課題2】学習に時間がかかる"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• GPU環境でも数時間～数日必要"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 1から構築では100文で約10時間以上"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "【アプローチ】特徴量（音素カバレッジ）に着目"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 音素カバレッジを最大化することで、データ量を削減"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Fine-tuningにより学習時間を短縮"
    p.font.size = Pt(18)
    p.level = 1


def add_phoneme_analysis_results(prs, phoneme_data):
    """音素解析結果スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "結果：音素解析"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "音素インベントリ"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = f"• 100コーパス中、合計{phoneme_data['phoneme_inventory_size']}音素（pau含む）を確認"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 最も多い音素: o（987回）"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 最も少ない音素: ny（29回）"
    p.font.size = Pt(18)
    p.level = 1
    
    # 音素頻度のTop10を表示
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "音素頻度 Top 10"
    p.font.size = Pt(20)
    p.font.bold = True
    p.level = 0
    
    freq = phoneme_data['phoneme_frequency']
    sorted_freq = sorted(freq.items(), key=lambda x: x[1], reverse=True)[:10]
    for i, (phoneme, count) in enumerate(sorted_freq):
        p = tf.add_paragraph()
        p.text = f"{i+1}. {phoneme}: {count}回"
        p.font.size = Pt(16)
        p.level = 1


def add_coverage_results(prs, corpus_data):
    """全音素カバーコーパス結果スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "結果：全音素(37音)カバーコーパス"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    e2_data = corpus_data["E2_4sent_37phonemes"]
    
    p = tf.paragraphs[0]
    p.text = "貪欲法による4文での全音素カバー"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 1コーパス目で84%"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 2コーパス目で92%"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 3コーパス目で97%"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 4コーパスで100%カバー"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "データ量削減効果"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 4文で全37音素をカバーできることを確認"
    p.font.size = Pt(18)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• データ量を96%削減（100文→4文）しながら全音素をカバー"
    p.font.size = Pt(18)
    p.font.bold = True
    p.level = 1


def add_finetuning_time_comparison(prs):
    """Fine-tuning時間比較スライドを追加"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "結果：Fine-tuning時間比較"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    # 表を作成
    rows = 5
    cols = 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2.5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # ヘッダー行
    table.cell(0, 0).text = "条件"
    table.cell(0, 1).text = "文数"
    table.cell(0, 2).text = "学習時間（10エポック）"
    
    # データ行（PDCAチェックリストから取得した時間）
    table.cell(1, 0).text = "80文コーパス"
    table.cell(1, 1).text = "80"
    table.cell(1, 2).text = "約37分40秒"
    
    table.cell(2, 0).text = "37音素4文"
    table.cell(2, 1).text = "4"
    table.cell(2, 2).text = "約45分"
    
    table.cell(3, 0).text = "ランダム4文"
    table.cell(3, 1).text = "4"
    table.cell(3, 2).text = "約37分46秒"
    
    table.cell(4, 0).text = "上位10文"
    table.cell(4, 1).text = "10"
    table.cell(4, 2).text = "約56分9秒"
    
    # フォントサイズ設定
    for row in range(rows):
        for col in range(cols):
            cell = table.cell(row, col)
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            if row == 0:
                cell.text_frame.paragraphs[0].font.bold = True
    
    # 比較テキストを追加
    left_text = Inches(1)
    top_text = Inches(4.8)
    width_text = Inches(8)
    height_text = Inches(1.5)
    
    text_box = slide.shapes.add_textbox(left_text, top_text, width_text, height_text)
    tf = text_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "1から構築との比較"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• 1から構築（100文で約10時間以上）と比較して大幅に短縮"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "• Fine-tuningにより学習時間を約93%削減"
    p.font.size = Pt(16)
    p.font.bold = True


def add_audio_slide(prs):
    """音声合成結果スライドを追加（主観評価用）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title = slide.shapes.title
    title.text = "結果：音声合成（主観評価用）"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "各条件の合成音声を比較"
    p.font.size = Pt(24)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "テキスト: 「スマートフォンから、フィーチャーフォンまで、マルチデバイスに対応。」"
    p.font.size = Pt(16)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    conditions = [
        ("ベースライン（事前学習モデル）", "baseline"),
        ("80文コーパス", "train_80sent"),
        ("37音素4文", "train_4sent_37phonemes"),
        ("ランダム4文", "train_4sent_random"),
        ("上位10文", "train_10sent_top")
    ]
    
    for condition_name, condition_dir in conditions:
        p = tf.add_paragraph()
        p.text = f"• {condition_name}"
        p.font.size = Pt(18)
        p.level = 1
        
        # 音声ファイルのパスを表示（実際の埋め込みは手動で行う）
        audio_path = PROJECT_ROOT / "outputs" / "audio" / condition_dir / "VOICEACTRESS100_010.wav"
        if audio_path.exists():
            p = tf.add_paragraph()
            p.text = f"  ファイル: {condition_dir}/VOICEACTRESS100_010.wav"
            p.font.size = Pt(14)
            p.font.italic = True
            p.level = 2
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "※ 音声ファイルは手動でPowerPointに埋め込んでください"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = RGBColor(150, 150, 150)
    
    p = tf.add_paragraph()
    p.text = "聴衆に主観評価を依頼（自然さ、話者類似度など）"
    p.font.size = Pt(16)
    p.font.bold = True


def main_final():
    """最終発表資料用のメイン関数"""
    print("最終発表資料を生成中...")
    
    # データの読み込み
    df = load_evaluation_data()
    phoneme_data = load_phoneme_data()
    corpus_data = load_corpus_selection()
    
    # プレゼンテーションオブジェクトの作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # スライドを追加
    print("  タイトルスライドを追加...")
    add_title_slide_final(prs)
    
    print("  目次スライドを追加...")
    add_table_of_contents(prs)
    
    print("  背景スライドを追加...")
    add_background_challenges(prs)
    
    print("  目的スライドを追加...")
    add_objective_slide_final(prs)
    
    print("  方法スライドを追加...")
    add_experiment_design_slide(prs)  # 方法を含む既存の関数を使用
    
    print("  音素解析結果スライドを追加...")
    add_phoneme_analysis_results(prs, phoneme_data)
    
    print("  全音素カバーコーパス結果スライドを追加...")
    add_coverage_results(prs, corpus_data)
    
    print("  Fine-tuning時間比較スライドを追加...")
    add_finetuning_time_comparison(prs)
    
    print("  評価指標スライドを追加...")
    add_evaluation_metrics_slide(prs)
    
    print("  MCD結果スライドを追加...")
    add_mcd_results_slide(prs, df)
    
    print("  log-F0 RMSE結果スライドを追加...")
    add_f0_results_slide(prs, df)
    
    print("  可視化スライドを追加...")
    add_visualization_slide(prs)
    
    print("  音声合成結果スライドを追加...")
    add_audio_slide(prs)
    
    print("  考察スライド1を追加...")
    add_discussion_slide_1(prs)
    
    print("  考察スライド2を追加...")
    add_discussion_slide_2(prs)
    
    print("  今後の課題スライドを追加...")
    add_future_work_slide(prs)
    
    # ファイルを保存
    output_file = OUTPUT_DIR / "slides_final.pptx"
    prs.save(str(output_file))
    
    print(f"\n最終発表資料を生成しました: {output_file}")
    print(f"総スライド数: {len(prs.slides)}枚")
    print("\n注意: 音声ファイルは手動でPowerPointに埋め込んでください")


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "--final":
        main_final()
    else:
        main()
